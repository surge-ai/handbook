import os
import pathlib
import shutil
import tempfile
import unittest
from unittest import mock

import pytest

from core.tools import sandbox
from core.tools.prepare_grading_context import (
    _extract_docx,
    _extract_pptx,
    _extract_text,
    _extract_xlsx,
    _format_evidence_as_xml,
    _get_extension,
    _is_excluded,
    _walk_dir,
    prepareGradingContext,
)


class TestHelpers(unittest.TestCase):
    def test_get_extension(self):
        assert _get_extension("/foo/bar.txt") == "txt"
        assert _get_extension("/foo/bar.CSV") == "csv"
        assert _get_extension("/foo/bar") == ""
        # Dotfiles like .hidden have no extension per os.path.splitext
        assert _get_extension("/foo/.hidden") == ""

    def test_is_excluded(self):
        excl = {"__pycache__", ".git"}
        assert _is_excluded("__pycache__/foo.pyc", excl)
        assert _is_excluded("src/__pycache__/foo.pyc", excl)
        assert _is_excluded(".git/config", excl)
        assert not _is_excluded("src/main.py", excl)

    def test_walk_dir_sorted(self):
        with tempfile.TemporaryDirectory() as d:
            os.makedirs(os.path.join(d, "b"))
            open(os.path.join(d, "b", "z.txt"), "w").close()
            open(os.path.join(d, "b", "a.txt"), "w").close()
            open(os.path.join(d, "c.txt"), "w").close()
            results = _walk_dir(d)
            names = [os.path.basename(r) for r in results]
            assert names == ["a.txt", "z.txt", "c.txt"] or results == sorted(results)


class TestExtractText(unittest.TestCase):
    # _extract_text now takes (raw_bytes, file_size, max_bytes) — the actual file
    # read happens earlier, as the sandbox user, in prepareGradingContext.
    def test_extract_text_normal(self):
        content, method, truncated = _extract_text(b"hello world", 11, 50_000)
        assert content == "hello world"
        assert method == "read"
        assert not truncated

    def test_extract_text_truncated(self):
        raw = b"a" * 1000
        content, _method, truncated = _extract_text(raw, len(raw), 100)
        assert truncated
        assert content is not None
        assert "[... truncated at 100 bytes" in content
        assert "total size: 1,000 bytes]" in content

    def test_extract_text_decodes_non_utf8_lossily(self):
        # Invalid UTF-8 bytes are replaced, never raise.
        content, _method, _truncated = _extract_text(b"\xff\xfe bad bytes", 12, 50_000)
        assert content is not None


class TestFormatEvidenceAsXml(unittest.TestCase):
    def test_empty_evidence(self):
        assert _format_evidence_as_xml([]) == ""

    def test_single_file_with_content(self):
        evidence = [
            {
                "path": "/tmp/test.txt",
                "extension": "txt",
                "size_bytes": 11,
                "content": "hello world",
                "truncated": False,
            }
        ]
        xml = _format_evidence_as_xml(evidence, "/tmp")
        assert '<workspace_files directory="/tmp">' in xml
        assert 'path="/tmp/test.txt"' in xml
        assert 'type="txt"' in xml
        assert 'size_bytes="11"' in xml
        assert "hello world" in xml
        assert "</workspace_files>" in xml

    def test_null_content_self_closing(self):
        evidence = [
            {
                "path": "/tmp/test.bin",
                "extension": "bin",
                "size_bytes": 100,
                "content": None,
                "truncated": False,
            }
        ]
        xml = _format_evidence_as_xml(evidence)
        assert "/>" in xml
        assert "</file>" not in xml

    def test_special_extension_extraction_method(self):
        evidence = [
            {
                "path": "/tmp/test.pdf",
                "extension": "pdf",
                "size_bytes": 5000,
                "content": "pdf text",
                "extraction_method": "pypdf",
                "truncated": False,
            }
        ]
        xml = _format_evidence_as_xml(evidence)
        assert 'extraction_method="pypdf"' in xml

    def test_truncated_attribute(self):
        evidence = [
            {
                "path": "/tmp/big.txt",
                "extension": "txt",
                "size_bytes": 100000,
                "content": "partial...",
                "truncated": True,
            }
        ]
        xml = _format_evidence_as_xml(evidence)
        assert 'truncated="true"' in xml


class TestPrepareGradingContext(unittest.IsolatedAsyncioTestCase):
    def setUp(self):
        self.temp_dir = tempfile.mkdtemp()
        self.original_workdir = sandbox.WORKDIR
        # prepareGradingContext now confines reads to the sandbox workdir, so
        # point WORKDIR at the dir under test (mirrors production, where grading
        # collects the agent's deliverables from WORKDIR).
        sandbox.WORKDIR = self.temp_dir

    def tearDown(self):
        shutil.rmtree(self.temp_dir, ignore_errors=True)
        sandbox.WORKDIR = self.original_workdir

    def _write(self, rel_path: str, content: str = "test") -> str:
        full = os.path.join(self.temp_dir, rel_path)
        os.makedirs(os.path.dirname(full), exist_ok=True)
        with open(full, "w") as f:
            f.write(content)
        return full

    async def test_returns_xml_prepended_to_final_output(self):
        self._write("hello.txt", "hello world")
        result = await prepareGradingContext(final_output="agent response", directory=self.temp_dir)
        assert result.startswith("<workspace_files")
        assert result.endswith("agent response")
        assert "\n\nagent response" in result

    async def test_no_files_returns_just_final_output(self):
        result = await prepareGradingContext(final_output="agent response", directory=self.temp_dir)
        assert result == "agent response"

    async def test_nonexistent_directory_returns_final_output(self):
        result = await prepareGradingContext(final_output="agent response", directory="/nonexistent/path")
        assert result == "agent response"

    async def test_default_workdir(self):
        workdir = tempfile.mkdtemp()
        sandbox.WORKDIR = workdir
        with open(os.path.join(workdir, "file.txt"), "w") as f:
            f.write("test")
        result = await prepareGradingContext(final_output="output")
        assert "<workspace_files" in result
        assert result.endswith("output")
        shutil.rmtree(workdir)

    async def test_extension_filter(self):
        self._write("hello.txt", "text")
        self._write("hello.py", "print('hi')")
        self._write("hello.csv", "a,b")
        result = await prepareGradingContext(
            final_output="output",
            directory=self.temp_dir,
            extensions=["txt", "csv"],
        )
        assert 'type="txt"' in result
        assert 'type="csv"' in result
        assert 'type="py"' not in result

    async def test_unsupported_extension_raises(self):
        with pytest.raises(ValueError) as ctx:
            await prepareGradingContext(
                final_output="output",
                directory=self.temp_dir,
                extensions=["exe"],
            )
        assert "Unsupported" in str(ctx.value)

    async def test_exclude_patterns(self):
        self._write("src/main.py", "code")
        self._write("__pycache__/cache.py", "cached")
        self._write(".git/config", "git stuff")
        result = await prepareGradingContext(final_output="output", directory=self.temp_dir)
        assert "main.py" in result
        assert "cache.py" not in result
        assert "config" not in result

    async def test_max_files_limit(self):
        for i in range(10):
            self._write(f"file{i:02d}.txt", f"content {i}")
        result = await prepareGradingContext(final_output="output", directory=self.temp_dir, max_files=3)
        assert result.count("<file ") == 3

    async def test_max_files_overflow_is_reported_not_silent(self):
        # The walk is flat + alphabetical, so a low cap can drop later files.
        # Those must be surfaced (count + names), not silently dropped, so a
        # rubric does not wrongly conclude an existing file is missing.
        for i in range(5):
            self._write(f"file{i:02d}.txt", f"content {i}")
        result = await prepareGradingContext(final_output="output", directory=self.temp_dir, max_files=2)
        assert result.count("<file ") == 2
        assert 'truncated_to_max_files="true"' in result
        assert 'dropped_file_count="3"' in result
        assert "<files_omitted_due_to_max_files" in result
        # The dropped files are the alphabetically-later ones.
        assert "file04.txt" in result
        assert result.count("<omitted_file ") == 3

    async def test_no_overflow_has_no_truncation_marker(self):
        for i in range(3):
            self._write(f"file{i:02d}.txt", f"content {i}")
        result = await prepareGradingContext(final_output="output", directory=self.temp_dir, max_files=10)
        assert "truncated_to_max_files" not in result
        assert "files_omitted_due_to_max_files" not in result

    async def test_overflow_marker_emitted_even_when_no_files_fit(self):
        # max_files=0: nothing is included, but the dir is non-empty. The grader
        # must still see that files exist (and were omitted), not a bare output.
        self._write("only.txt", "content")
        result = await prepareGradingContext(final_output="output", directory=self.temp_dir, max_files=0)
        assert result.startswith("<workspace_files")
        assert 'dropped_file_count="1"' in result
        assert "only.txt" in result

    async def test_unsupported_files_skipped(self):
        self._write("image.png", "binary")
        self._write("doc.txt", "text")
        result = await prepareGradingContext(final_output="output", directory=self.temp_dir)
        assert 'type="txt"' in result
        # Check for the unsupported file's marker tokens specifically — the
        # tempdir path itself can contain "png" as a random substring.
        assert "image.png" not in result
        assert 'type="png"' not in result

    async def test_oversized_special_file_recorded_not_extracted(self):
        # A special file over the raw-read cap must be surfaced as evidence (so
        # the rubric knows it exists) but NOT extracted — extracting truncated
        # PDF/Office bytes would corrupt, and reading it whole could OOM grading.
        self._write("big.pdf", "x" * 100)
        with mock.patch("core.tools.prepare_grading_context.MAX_SPECIAL_FILE_BYTES", 10):
            result = await prepareGradingContext(final_output="output", directory=self.temp_dir)
        assert 'type="pdf"' in result
        assert 'size_bytes="100"' in result
        assert "exceeds" in result and "grading read cap" in result
        assert 'truncated="true"' in result
        # content is None -> self-closing tag, and the raw bytes never appear.
        assert "<big.pdf" not in result

    async def test_special_file_under_cap_is_extracted(self):
        # A small docx under the cap still flows through normal extraction — the
        # cap must not block legitimate files.
        from docx import Document

        path = os.path.join(self.temp_dir, "report.docx")
        doc = Document()
        doc.add_paragraph("under-cap evidence")
        doc.save(path)
        result = await prepareGradingContext(final_output="output", directory=self.temp_dir)
        assert 'extraction_method="python-docx"' in result
        assert "under-cap evidence" in result

    async def test_returns_plain_string(self):
        self._write("test.txt", "content")
        result = await prepareGradingContext(final_output="my output", directory=self.temp_dir)
        assert isinstance(result, str)
        # Should NOT be JSON
        assert not result.startswith("{")
        assert not result.startswith("[")

    async def test_empty_directory_returns_final_output(self):
        result = await prepareGradingContext(final_output="my output", directory=self.temp_dir)
        assert result == "my output"


class TestExtractDocx(unittest.TestCase):
    def _make_docx(self, path: str) -> None:
        from docx import Document

        doc = Document()
        doc.add_paragraph("Quarterly summary: revenue up 20%.")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Revenue"
        table.cell(1, 1).text = "1000"
        doc.save(path)

    def test_extracts_paragraphs(self):
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "report.docx")
            self._make_docx(path)
            content, method, truncated = _extract_docx(pathlib.Path(path).read_bytes(), 50_000)
            assert method == "python-docx"
            assert not truncated
            assert content is not None
            assert "revenue up 20%" in content

    def test_includes_table_cells(self):
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "report.docx")
            self._make_docx(path)
            content, _method, _truncated = _extract_docx(pathlib.Path(path).read_bytes(), 50_000)
            # Regression: table cell text used to be dropped (paragraphs only).
            assert content is not None
            assert "Metric" in content
            assert "Revenue" in content
            assert "1000" in content

    def test_includes_nested_table_cells(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "nested.docx")
            doc = Document()
            outer = doc.add_table(rows=1, cols=1)
            cell = outer.cell(0, 0)
            cell.paragraphs[0].text = "OuterCell"
            inner = cell.add_table(rows=1, cols=2)
            inner.cell(0, 0).text = "InnerA"
            inner.cell(0, 1).text = "InnerB"
            doc.save(path)

            content, _method, _truncated = _extract_docx(pathlib.Path(path).read_bytes(), 50_000)
            assert content is not None
            # cell.text flattens only direct paragraphs; nested tables must recurse.
            assert "OuterCell" in content
            assert "InnerA" in content
            assert "InnerB" in content

    def test_missing_file_returns_none(self):
        content, method, truncated = _extract_docx(b"not a real docx", 50_000)
        assert content is None
        assert method == "python-docx"
        assert not truncated


class TestExtractXlsx(unittest.TestCase):
    def _make_xlsx(self, path: str) -> None:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Value"])
        ws.append(["foo", 42])
        wb.save(path)

    def test_extracts_rows(self):
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "sheet.xlsx")
            self._make_xlsx(path)
            content, method, truncated = _extract_xlsx(pathlib.Path(path).read_bytes(), 50_000)
            assert method == "openpyxl"
            assert not truncated
            assert content is not None
            assert "--- Sheet: Data ---" in content
            assert "Name\tValue" in content
            assert "foo\t42" in content

    def test_missing_file_returns_none(self):
        content, method, truncated = _extract_xlsx(b"not a real xlsx", 50_000)
        assert content is None
        assert method == "openpyxl"
        assert not truncated

    def test_every_sheet_appears_even_when_earlier_sheet_truncated(self):
        # Regression: a big first sheet used to drain the whole budget, so later
        # tabs never appeared. The budget is now per-sheet.
        from openpyxl import Workbook

        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "big.xlsx")
            wb = Workbook()
            ws1 = wb.active
            ws1.title = "TB Current"
            for i in range(2000):
                ws1.append([f"row{i}", "Trial Balance as of 12/31/2025", i, i * 2, "padding text"])
            ws2 = wb.create_sheet("TB Prior")
            ws2.append(["Trial Balance as of 11/30/2025"])
            ws2.append(["Prior period data", 123])
            wb.save(path)

            content, _method, truncated = _extract_xlsx(pathlib.Path(path).read_bytes(), 2000)
            assert truncated
            assert content is not None
            assert "--- Sheet: TB Current ---" in content
            assert "--- Sheet: TB Prior ---" in content
            # The starved later sheet's data survives, and truncation is attributed.
            assert "11/30/2025" in content
            assert "Prior period data\t123" in content
            assert "'TB Current' truncated" in content

    def test_empty_trailing_rows_do_not_consume_budget(self):
        # openpyxl read-only dimensions can be inflated; fully-empty rows must be
        # skipped so they don't burn the budget producing only blank tab output.
        from openpyxl import Workbook

        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "sparse.xlsx")
            wb = Workbook()
            ws = wb.active
            ws.title = "Sparse"
            ws.append(["header", "value"])
            ws.append(["real", 1])
            # Force a large inflated dimension with empty cells far below.
            ws.cell(row=5000, column=1, value=None)
            wb.save(path)

            content, _method, truncated = _extract_xlsx(pathlib.Path(path).read_bytes(), 50_000)
            assert content is not None
            assert not truncated
            assert "header\tvalue" in content
            assert "real\t1" in content
            # No run of empty tab-joined rows should appear.
            assert "\t\t\t" not in content


class TestExtractPptx(unittest.TestCase):
    def _make_pptx(self, path: str) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[5])
        s1.shapes.title.text = "Helix Pricing SLT Deck"
        # Slide 2: a textbox, a table, and speaker notes.
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Recommended price: $499"
        table = s2.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4), Inches(1)).table
        table.cell(0, 0).text = "Tier"
        table.cell(0, 1).text = "Price"
        table.cell(1, 0).text = "Pro"
        table.cell(1, 1).text = "499"
        s2.notes_slide.notes_text_frame.text = "Emphasize ROI in the meeting"
        prs.save(path)

    def test_extracts_slides_tables_and_notes(self):
        with tempfile.TemporaryDirectory() as d:
            path = os.path.join(d, "deck.pptx")
            self._make_pptx(path)
            content, method, truncated = _extract_pptx(pathlib.Path(path).read_bytes(), 50_000)
            assert method == "python-pptx"
            assert not truncated
            assert content is not None
            assert "--- Slide 1 ---" in content
            assert "--- Slide 2 ---" in content
            assert "Helix Pricing SLT Deck" in content
            assert "Recommended price: $499" in content
            assert "Tier\tPrice" in content
            assert "Pro\t499" in content
            assert "[Notes] Emphasize ROI in the meeting" in content

    def test_missing_file_returns_none(self):
        content, method, truncated = _extract_pptx(b"not a real pptx", 50_000)
        assert content is None
        assert method == "python-pptx"
        assert not truncated


class TestPrepareGradingContextBinary(unittest.IsolatedAsyncioTestCase):
    async def test_docx_and_xlsx_evidence_assembled(self):
        from docx import Document
        from openpyxl import Workbook

        original_workdir = sandbox.WORKDIR
        with tempfile.TemporaryDirectory() as d:
            sandbox.WORKDIR = d  # reads are confined to WORKDIR
            self.addCleanup(setattr, sandbox, "WORKDIR", original_workdir)
            doc = Document()
            doc.add_paragraph("agent wrote this in a docx")
            doc.save(os.path.join(d, "out.docx"))

            wb = Workbook()
            wb.active.append(["col", "val"])
            wb.save(os.path.join(d, "out.xlsx"))

            from pptx import Presentation

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "agent wrote this in a pptx"
            prs.save(os.path.join(d, "out.pptx"))

            result = await prepareGradingContext(final_output="FINAL", directory=d)
            assert 'type="docx"' in result
            assert 'extraction_method="python-docx"' in result
            assert "agent wrote this in a docx" in result
            assert 'type="xlsx"' in result
            assert 'extraction_method="openpyxl"' in result
            # Regression: a .pptx deliverable used to be silently skipped.
            assert 'type="pptx"' in result
            assert 'extraction_method="python-pptx"' in result
            assert "agent wrote this in a pptx" in result
            assert result.endswith("FINAL")


if __name__ == "__main__":
    unittest.main()
