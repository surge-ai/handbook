"""Prepare grading context for LLM-based rubric evaluation.

Collects file evidence from the sandbox, formats as XML, and assembles
the full text_for_grading string by prepending file evidence to the
agent's final output.
"""

import asyncio
import io
import logging
import os

from core.tools import sandbox

logger = logging.getLogger(__name__)

# Default cap on the number of files included as evidence. Sized to comfortably
# cover a task's read-only input documents PLUS the agent's deliverables: the
# walk is flat + alphabetical, so input files can otherwise consume every slot
# and silently drop the agent-created files the rubric actually grades. When the
# cap IS hit, the dropped files are surfaced both in a log line and as a
# <files_omitted_due_to_max_files> marker inside the <workspace_files> XML, so an
# empty or partial listing is never mistaken for "the agent created nothing."
DEFAULT_MAX_FILES = 100

# Per-file content budget that ends up in the grading prompt. The ceiling on the
# whole evidence block is DEFAULT_MAX_FILES * this, so the constraint isn't memory
# but the grader's context window and cost — every byte here is fed to the rubric
# LLM. 250 KB gives large multi-tab spreadsheets and multi-slide decks room to
# show every tab/slide (the per-sheet/per-slide budgeting below distributes it),
# while keeping the worst case bounded. Realistic tasks have a handful of
# substantive files, so actual prompts stay far below the 100-file ceiling.
DEFAULT_MAX_CONTENT_BYTES = 250_000

TEXT_EXTENSIONS = {
    "txt",
    "csv",
    "json",
    "jsonl",
    "py",
    "js",
    "ts",
    "md",
    "html",
    "xml",
    "yaml",
    "yml",
    "toml",
    "cfg",
    "ini",
    "log",
    "sh",
    "bash",
    "sql",
    "r",
    "rb",
    "java",
    "c",
    "cpp",
    "h",
    "hpp",
    "css",
    "scss",
    "tex",
    "rst",
}

SPECIAL_EXTENSIONS = {"pdf", "docx", "xlsx", "pptx"}

ALL_SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | SPECIAL_EXTENSIONS

# Special formats (PDF/Office) can't be truncated without corrupting the
# container, so extraction needs the whole file in memory — but an unbounded
# read lets one giant agent-generated file OOM grading before the rubric runs.
# Cap the raw read: files over this size are recorded as evidence (so the rubric
# still knows they exist) but not extracted. Plain text is already bounded by
# max_content_bytes, so it needs no separate cap.
MAX_SPECIAL_FILE_BYTES = 50 * 1024 * 1024

DEFAULT_EXCLUDE_PATTERNS = {
    "__pycache__",
    ".git",
    ".venv",
    "node_modules",
    ".mypy_cache",
    ".pytest_cache",
    ".ruff_cache",
}


def _is_excluded(relative_path: str, exclude_patterns: set[str]) -> bool:
    parts = relative_path.split(os.sep)
    return any(part in exclude_patterns for part in parts)


def _within_workdir(path: str, workdir_root: str) -> bool:
    """True if ``path`` resolves to ``workdir_root`` or somewhere beneath it.

    ``realpath`` collapses ``..`` and follows symlinks, so a symlink inside the
    workdir pointing at ``/app`` resolves outside ``workdir_root`` and is
    rejected. The core server runs as root, so this is what stops
    ``prepareGradingContext`` from being turned into a read oracle for the
    locked-down ``/app`` tree.
    """
    canonical = os.path.realpath(path)
    return canonical == workdir_root or canonical.startswith(workdir_root + os.sep)


def _get_extension(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    return ext.lstrip(".")


def _extract_text(raw: bytes, file_size: int, max_bytes: int) -> tuple[str | None, str, bool]:
    try:
        truncated = file_size > max_bytes
        content = raw[:max_bytes].decode("utf-8", errors="replace")
        if truncated:
            content += f"\n\n[... truncated at {max_bytes:,} bytes, total size: {file_size:,} bytes]"
        return content, "read", truncated
    except Exception:
        return None, "read", False


def _extract_pdf(raw: bytes, max_bytes: int) -> tuple[str | None, str, bool]:
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(raw))
        parts: list[str] = []
        total_bytes = 0
        truncated = False
        for page in reader.pages:
            text = page.extract_text() or ""
            text_bytes = len(text.encode("utf-8"))
            if total_bytes + text_bytes > max_bytes:
                remaining = max_bytes - total_bytes
                parts.append(text[:remaining])
                truncated = True
                break
            parts.append(text)
            total_bytes += text_bytes
        content = "\n".join(parts)
        if truncated:
            content += f"\n\n[... truncated at {max_bytes:,} bytes]"
        return content, "pypdf", truncated
    except Exception:
        return None, "pypdf", False


def _docx_table_lines(table) -> list[str]:
    """Return one tab-separated line per row, recursing into nested tables.

    A cell can hold both paragraphs and further tables, so we join the cell's
    own paragraph text and then descend into any tables it nests — `cell.text`
    alone flattens only the direct paragraphs and silently drops nested tables.
    """
    lines: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        nested: list[str] = []
        for cell in row.cells:
            cells.append("\n".join(p.text for p in cell.paragraphs))
            for inner in cell.tables:
                nested.extend(_docx_table_lines(inner))
        lines.append("\t".join(cells))
        lines.extend(nested)
    return lines


def _extract_docx(raw: bytes, max_bytes: int) -> tuple[str | None, str, bool]:
    try:
        from docx import Document

        doc = Document(io.BytesIO(raw))
        lines = [p.text for p in doc.paragraphs]
        # Paragraphs alone miss text inside tables, which agents routinely use
        # for structured output (reports, spreadsheets exported as docx, etc.).
        for table in doc.tables:
            lines.extend(_docx_table_lines(table))
        text = "\n".join(lines)
        text_bytes = len(text.encode("utf-8"))
        truncated = text_bytes > max_bytes
        if truncated:
            content = text[:max_bytes]
            content += f"\n\n[... truncated at {max_bytes:,} bytes]"
        else:
            content = text
        return content, "python-docx", truncated
    except Exception:
        return None, "python-docx", False


def _extract_xlsx(raw: bytes, max_bytes: int) -> tuple[str | None, str, bool]:
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts: list[str] = []
        total_bytes = 0
        truncated = False
        # Budget the bytes PER SHEET instead of letting the first sheet drain the
        # whole allowance. A multi-tab workbook used to truncate mid-first-sheet
        # and never emit the later tabs at all — so a grader couldn't even tell
        # those sheets existed. We give each not-yet-visited sheet an equal slice
        # of the remaining budget, recomputed after every sheet so unused bytes
        # from small/empty sheets roll forward to later ones.
        sheet_names = wb.sheetnames
        for index, sheet_name in enumerate(sheet_names):
            ws = wb[sheet_name]
            # Always emit the header so every tab is visible, even one whose rows
            # get fully truncated away.
            header = f"--- Sheet: {sheet_name} ---"
            parts.append(header)
            total_bytes += len(header.encode("utf-8")) + 1

            remaining_sheets = len(sheet_names) - index
            sheet_budget = max(0, (max_bytes - total_bytes) // remaining_sheets)
            sheet_used = 0
            sheet_truncated = False
            for row in ws.iter_rows(values_only=True):
                line = "\t".join(str(cell) if cell is not None else "" for cell in row)
                # Skip fully-empty rows: openpyxl's read-only dimensions can be
                # inflated, so a sheet may report thousands of blank trailing rows
                # that would otherwise burn the budget producing only tabs.
                if not line.strip():
                    continue
                line_bytes = len(line.encode("utf-8")) + 1  # +1 for newline
                if sheet_used + line_bytes > sheet_budget:
                    sheet_truncated = True
                    truncated = True
                    break
                parts.append(line)
                sheet_used += line_bytes
                total_bytes += line_bytes
            if sheet_truncated:
                parts.append(f"[... '{sheet_name}' truncated at {sheet_budget:,} bytes for this sheet]")
        wb.close()
        content = "\n".join(parts)
        if truncated:
            content += f"\n\n[... some sheets truncated; per-sheet budget of ~{max_bytes:,} total bytes reached]"
        return content, "openpyxl", truncated
    except Exception:
        return None, "openpyxl", False


def _pptx_shape_lines(shapes) -> list[str]:
    """Return text lines for a collection of shapes, recursing into groups.

    Pulls text frames and table cells, descending into grouped shapes.
    """
    lines: list[str] = []
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            lines.extend(_pptx_shape_lines(shape.shapes))
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text
            if text.strip():
                lines.append(text)
        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text for cell in row.cells]
                lines.append("\t".join(cells))
    return lines


def _extract_pptx(raw: bytes, max_bytes: int) -> tuple[str | None, str, bool]:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw))
        parts: list[str] = []
        total_bytes = 0
        truncated = False
        for index, slide in enumerate(prs.slides, start=1):
            # Per-slide marker so the grader sees the slide count, which is often
            # itself part of the rubric.
            header = f"--- Slide {index} ---"
            lines = _pptx_shape_lines(slide.shapes)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    lines.append(f"[Notes] {notes}")
            block = "\n".join([header, *lines])
            block_bytes = len(block.encode("utf-8")) + 1
            if total_bytes + block_bytes > max_bytes:
                remaining = max_bytes - total_bytes
                if remaining > 0:
                    parts.append(block[:remaining])
                truncated = True
                break
            parts.append(block)
            total_bytes += block_bytes
        content = "\n".join(parts)
        if truncated:
            content += f"\n\n[... truncated at {max_bytes:,} bytes]"
        return content, "python-pptx", truncated
    except Exception:
        return None, "python-pptx", False


def _walk_dir(directory: str) -> list[str]:
    results: list[str] = []
    for root, _dirs, files in os.walk(directory):
        for name in files:
            results.append(os.path.join(root, name))
    results.sort()
    return results


def _format_evidence_as_xml(
    evidence: list[dict],
    directory: str | None = None,
    dropped_paths: list[str] | None = None,
) -> str:
    # Emit the container even with no included files when some were dropped by
    # the cap, so the grader sees the truncation notice rather than a bare
    # final_output. Only an empty listing AND nothing dropped means "no files."
    if not evidence and not dropped_paths:
        return ""

    dir_attr = f' directory="{directory}"' if directory else ""
    truncated_attr = ""
    if dropped_paths:
        truncated_attr = f' truncated_to_max_files="true" dropped_file_count="{len(dropped_paths)}"'
    parts: list[str] = [f"<workspace_files{dir_attr}{truncated_attr}>"]

    for item in evidence:
        attrs = f'path="{item["path"]}" type="{item["extension"]}" size_bytes="{item["size_bytes"]}"'
        if item.get("extraction_method"):
            attrs += f' extraction_method="{item["extraction_method"]}"'
        if item.get("truncated"):
            attrs += ' truncated="true"'

        if item["content"] is not None:
            parts.append(f"  <file {attrs}>")
            parts.append(item["content"])
            parts.append("  </file>")
        else:
            parts.append(f"  <file {attrs} />")

    if dropped_paths:
        # Surface the names of files that were walked but excluded by the cap.
        # An absent file here is genuinely absent; a file listed here exists but
        # its contents were not included, so a rubric must not conclude it is
        # missing from the workspace.
        parts.append(f'  <files_omitted_due_to_max_files count="{len(dropped_paths)}">')
        for path in dropped_paths:
            parts.append(f'    <omitted_file path="{path}" />')
        parts.append("  </files_omitted_due_to_max_files>")

    parts.append("</workspace_files>")
    return "\n".join(parts)


async def prepareGradingContext(
    final_output: str,
    directory: str | None = None,
    extensions: list[str] | None = None,
    exclude_patterns: list[str] | None = None,
    max_files: int = DEFAULT_MAX_FILES,
    max_content_bytes: int = DEFAULT_MAX_CONTENT_BYTES,
) -> str:
    """Prepare the text_for_grading string for rubric evaluation.

    Collects file evidence from the sandbox, formats as XML,
    and prepends to the agent's final_output.

    Args:
        final_output: The agent's final response text.
        directory: Root directory to search. Defaults to the sandbox directory.
        extensions: Only include files with these extensions (without dot).
        exclude_patterns: Directory/file name patterns to skip.
        max_files: Maximum number of files to include as evidence. Eligible files
            beyond this cap are omitted but reported (logged and listed in the
            XML under <files_omitted_due_to_max_files>) rather than dropped
            silently. Defaults to DEFAULT_MAX_FILES.
        max_content_bytes: Maximum bytes of content to read per file.

    Returns:
        The assembled text_for_grading string.
    """
    # Walks the sandbox tree and reads/parses files (text/PDF/docx) — all
    # blocking — so run the whole body in a worker thread to free the loop.
    return await asyncio.to_thread(
        _prepare_grading_context_sync,
        final_output,
        directory,
        extensions,
        exclude_patterns,
        max_files,
        max_content_bytes,
    )


def _prepare_grading_context_sync(
    final_output: str,
    directory: str | None = None,
    extensions: list[str] | None = None,
    exclude_patterns: list[str] | None = None,
    max_files: int = DEFAULT_MAX_FILES,
    max_content_bytes: int = DEFAULT_MAX_CONTENT_BYTES,
) -> str:
    if directory is None:
        directory = sandbox.WORKDIR

    if not os.path.isdir(directory):
        return final_output

    # Grading only ever needs the agent's deliverables, which live under the
    # sandbox workdir. The server runs as root, so confine the walk to the
    # workdir — anything outside it (a caller passing /app, or a symlink that
    # escapes) is refused so this can't exfiltrate the locked-down tree.
    workdir_root = os.path.realpath(sandbox.WORKDIR)
    if not _within_workdir(directory, workdir_root):
        logger.warning(
            "prepareGradingContext: directory %r is outside the sandbox workdir %r; refusing to read it",
            directory,
            sandbox.WORKDIR,
        )
        return final_output

    excl = set(exclude_patterns) if exclude_patterns else DEFAULT_EXCLUDE_PATTERNS

    ext_filter: set[str] | None = None
    if extensions:
        ext_filter = {e.lower().lstrip(".") for e in extensions}
        unsupported = ext_filter - ALL_SUPPORTED_EXTENSIONS
        if unsupported:
            raise ValueError(
                f"Unsupported file extension(s): {', '.join('.' + e for e in sorted(unsupported))}. "
                f"Supported extensions: {', '.join('.' + e for e in sorted(ALL_SUPPORTED_EXTENSIONS))}"
            )

    all_files = _walk_dir(directory)
    evidence: list[dict] = []
    # Files that passed every filter (would have been included) but fell outside
    # the max_files cap. Tracked so the cap is observable rather than silent: a
    # rubric seeing one of these must not conclude the file is absent.
    dropped_paths: list[str] = []

    for file_path in all_files:
        # Skip symlinks (and any other entry) whose real target escapes the
        # workdir — they'd otherwise let a root read reach outside the sandbox.
        if not _within_workdir(file_path, workdir_root):
            continue
        relative_path = os.path.relpath(file_path, directory)
        if _is_excluded(relative_path, excl):
            continue

        ext = _get_extension(file_path)
        if ext not in ALL_SUPPORTED_EXTENSIONS:
            continue
        if ext_filter and ext not in ext_filter:
            continue

        # Cap is checked AFTER filtering so it counts only eligible files, and
        # records the overflow instead of silently breaking out of the walk.
        if len(evidence) >= max_files:
            dropped_paths.append(file_path)
            continue

        # Read bytes AS THE SANDBOX USER. The server is root, so opening the file
        # in-process would both bypass /app's permissions and be a symlink-swap
        # TOCTOU oracle; reading as uid 1000 closes both. Office/PDF formats can't
        # be truncated (zip/PDF bytes corrupt) so extraction needs the whole file;
        # plain text is bounded by max_content_bytes.
        read_limit = MAX_SPECIAL_FILE_BYTES if ext in SPECIAL_EXTENSIONS else max_content_bytes

        # For special formats, probe the size first (limit=0 reads no bytes) so an
        # oversized file is recorded as evidence and skipped — without streaming
        # MAX_SPECIAL_FILE_BYTES through the pipe just to discard it, which would
        # let a few huge agent files waste GBs and risk OOMing grading.
        if ext in SPECIAL_EXTENSIONS:
            try:
                size_bytes, _h, _r, _s = sandbox.agent_read_window(file_path, offset=0, limit=0, sniff=0)
            except sandbox.AgentReadError:
                continue
            if size_bytes > MAX_SPECIAL_FILE_BYTES:
                evidence.append(
                    {
                        "path": file_path,
                        "extension": ext,
                        "size_bytes": size_bytes,
                        "content": None,
                        "truncated": True,
                        "extraction_method": (
                            f"skipped: {size_bytes:,} bytes exceeds the {MAX_SPECIAL_FILE_BYTES:,}-byte grading read cap"
                        ),
                    }
                )
                continue

        try:
            size_bytes, _header, raw, _start = sandbox.agent_read_window(file_path, offset=0, limit=read_limit, sniff=0)
        except sandbox.AgentReadError:
            continue

        if ext == "pdf":
            content, method, truncated = _extract_pdf(raw, max_content_bytes)
        elif ext == "docx":
            content, method, truncated = _extract_docx(raw, max_content_bytes)
        elif ext == "xlsx":
            content, method, truncated = _extract_xlsx(raw, max_content_bytes)
        elif ext == "pptx":
            content, method, truncated = _extract_pptx(raw, max_content_bytes)
        else:
            content, method, truncated = _extract_text(raw, size_bytes, max_content_bytes)

        item: dict = {
            "path": file_path,
            "extension": ext,
            "size_bytes": size_bytes,
            "content": content,
            "truncated": truncated,
        }
        if ext in SPECIAL_EXTENSIONS:
            item["extraction_method"] = method

        evidence.append(item)

    if dropped_paths:
        logger.warning(
            "prepareGradingContext hit max_files=%d in %s: included %d file(s), "
            "omitted %d eligible file(s) from grading evidence: %s",
            max_files,
            directory,
            len(evidence),
            len(dropped_paths),
            dropped_paths,
        )

    xml = _format_evidence_as_xml(evidence, directory, dropped_paths)

    if xml:
        return f"{xml}\n\n{final_output}"
    return final_output
